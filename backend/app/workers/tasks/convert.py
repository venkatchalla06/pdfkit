import os
import subprocess
import tempfile
import zipfile

import fitz
from PIL import Image

from app.workers.celery_app import celery_app
from app.workers.base_task import PDFBaseTask, SyncSession
from app.db.models.job import Job, JobStatus, ToolType
from app.services.storage import storage


def libreoffice_convert(input_path: str, output_dir: str, fmt: str = "pdf") -> str:
    """Convert using LibreOffice headless. Returns path of output file."""
    before = set(os.listdir(output_dir))

    # Each job gets its own LO user-profile so concurrent conversions don't
    # lock each other out ("source file could not be loaded" race condition).
    lo_profile = os.path.join(output_dir, ".lo_profile")
    os.makedirs(lo_profile, exist_ok=True)

    cmd = [
        "libreoffice",
        f"-env:UserInstallation=file://{lo_profile}",
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        "--convert-to", fmt,
        "--outdir", output_dir,
        input_path,
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=120)
    except FileNotFoundError:
        raise RuntimeError("LibreOffice is not installed. Install libreoffice-writer.")

    stderr = result.stderr.decode()
    stdout = result.stdout.decode()

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed (code {result.returncode}): {stderr[:400]}")

    # Detect the output file by directory diff (handles any name LibreOffice chose)
    after = set(os.listdir(output_dir))
    new_files = [f for f in (after - before) if f.lower().endswith(f".{fmt}")]
    if new_files:
        return os.path.join(output_dir, new_files[0])

    expected = os.path.join(output_dir, os.path.splitext(os.path.basename(input_path))[0] + f".{fmt}")
    if os.path.exists(expected):
        return expected

    raise RuntimeError(
        f"LibreOffice produced no output. stdout: {stdout[:200]} stderr: {stderr[:200]}"
    )


# ── PDF → JPG ──────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_jpg_task",
    max_retries=2, soft_time_limit=300,
)
def pdf_to_jpg_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    dpi = int(options.get("dpi", 150))

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)

        doc = fitz.open(input_path)
        total = len(doc)
        image_paths = []
        mat = fitz.Matrix(dpi / 72, dpi / 72)

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_path = os.path.join(tmpdir, f"page_{i + 1:04d}.jpg")
            pix.save(img_path)
            image_paths.append(img_path)
            self.update_job(job_id, progress=5 + int((i + 1) / total * 75))

        doc.close()

        if total == 1:
            output_key = f"results/{job_id}/page_1.jpg"
            storage.upload_from_temp(image_paths[0], output_key, content_type="image/jpeg")
            output_filename = "page_1.jpg"
        else:
            zip_path = os.path.join(tmpdir, "pages.zip")
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                for p in image_paths:
                    zf.write(p, os.path.basename(p))
            output_key = f"results/{job_id}/pages.zip"
            storage.upload_from_temp(zip_path, output_key, content_type="application/zip")
            output_filename = "pages.zip"

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": output_filename}
        session.commit()

    return output_key


# ── JPG → PDF ──────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.jpg_to_pdf_task",
    max_retries=2, soft_time_limit=120,
)
def jpg_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_keys = job.input_keys
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        doc = fitz.open()

        for i, key in enumerate(input_keys):
            img_path = os.path.join(tmpdir, f"img_{i}.jpg")
            storage.download_to_temp(key, img_path)

            # Use PIL to get dimensions
            with Image.open(img_path) as img:
                w, h = img.size

            page = doc.new_page(width=w, height=h)
            page.insert_image(page.rect, filename=img_path)
            self.update_job(job_id, progress=5 + int((i + 1) / len(input_keys) * 80))

        output_path = os.path.join(tmpdir, "images.pdf")
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()

        output_key = f"results/{job_id}/images.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "images.pdf"}
        session.commit()

    return output_key


# ── PDF → Word ─────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_word_task",
    max_retries=1, soft_time_limit=300,
)
def pdf_to_word_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        output_path = os.path.join(tmpdir, "converted.docx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        # pdf2docx handles PDF→DOCX far better than LibreOffice
        from pdf2docx import Converter
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/converted.docx"
        storage.upload_from_temp(
            output_path, output_key,
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.docx"}
        session.commit()

    return output_key


# ── Word → PDF ─────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.word_to_pdf_task",
    max_retries=1, soft_time_limit=300,
)
def word_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.docx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        output_path = libreoffice_convert(input_path, tmpdir, fmt="pdf")
        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/converted.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.pdf"}
        session.commit()

    return output_key


# ── PowerPoint → PDF ───────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pptx_to_pdf_task",
    max_retries=1, soft_time_limit=300,
)
def pptx_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pptx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        output_path = libreoffice_convert(input_path, tmpdir, fmt="pdf")
        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/presentation.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "presentation.pdf"}
        session.commit()

    return output_key


# ── Excel → PDF ────────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.xlsx_to_pdf_task",
    max_retries=1, soft_time_limit=300,
)
def xlsx_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.xlsx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        output_path = libreoffice_convert(input_path, tmpdir, fmt="pdf")
        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/spreadsheet.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "spreadsheet.pdf"}
        session.commit()

    return output_key


# ── HTML → PDF ─────────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.html_to_pdf_task",
    max_retries=1, soft_time_limit=180,
)
def html_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.html")
        storage.download_to_temp(input_key, input_path)

        self.update_job(job_id, progress=30)

        lo_output = libreoffice_convert(input_path, tmpdir, fmt="pdf")

        output_key = f"results/{job_id}/converted.pdf"
        storage.upload_from_temp(lo_output, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.pdf"}
        session.commit()

    return output_key


# ── PDF → PowerPoint ───────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_pptx_task",
    max_retries=1, soft_time_limit=300,
)
def pdf_to_pptx_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)

        from pptx import Presentation

        doc = fitz.open(input_path)
        total = len(doc)
        prs = Presentation()

        for i, page in enumerate(doc):
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_path = os.path.join(tmpdir, f"slide_{i:04d}.jpg")
            pix.save(img_path)

            slide_w = int(page.rect.width * 914400 / 72)
            slide_h = int(page.rect.height * 914400 / 72)
            prs.slide_width = slide_w
            prs.slide_height = slide_h

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)

            self.update_job(job_id, progress=5 + int((i + 1) / total * 85))

        doc.close()

        output_path = os.path.join(tmpdir, "converted.pptx")
        prs.save(output_path)

        output_key = f"results/{job_id}/converted.pptx"
        storage.upload_from_temp(
            output_path, output_key,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.pptx"}
        session.commit()

    return output_key
